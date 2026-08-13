"""Executive report generation (Excel + PowerPoint) from the live application analytics.

Every value comes from the same Fleet-scoped analytics the UI uses (overview, sla,
telemetry, trends, resilience) so exported figures reconcile with the dashboard. Nothing
is hard-coded; missing evidence stays "Insufficient"/"Not monitored". The visual design
follows the Medline executive template: a deep-blue (#003DA5) brand banner, white KPI
cards with large navy values and colour-coded status, an executive readout with status
chips, and a per-site availability bar chart. Read-only (no LogicMonitor write access).
"""
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.orm import Session

from . import overview, resilience, sla, telemetry, trends
from .config import get_settings
from .db import SlaDaily

# ---- Medline executive palette (hex for Excel, RGBColor for PowerPoint) ----
BRAND = "003DA5"   # Medline blue — banners / brand
NAVY = "0F172A"    # KPI values / primary headings
SLATE = "334155"   # labels / secondary text
ZEBRA = "F1F5F9"   # alternating rows / table header fill
CANVAS = "F8FAFC"  # slide canvas
LINE = "E2E8F0"    # 1px card outline
GREEN, GREEN_BG = "16A34A", "DCFCE7"
RED, RED_BG = "C00000", "FEE2E2"
AMBER, AMBER_BG = "CA8A04", "FEF9C3"
BLUE_BG = "DBEAFE"

_BRAND = RGBColor(0x00, 0x3D, 0xA5)
_NAVY = RGBColor(0x0F, 0x17, 0x2A)
_SLATE = RGBColor(0x33, 0x41, 0x55)
_GREY = RGBColor(0x64, 0x74, 0x8B)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LINE = RGBColor(0xE2, 0xE8, 0xF0)
_CANVAS = RGBColor(0xF8, 0xFA, 0xFC)
_GREEN = RGBColor(0x16, 0xA3, 0x4A)
_RED = RGBColor(0xC0, 0x00, 0x00)
_AMBER = RGBColor(0xCA, 0x8A, 0x04)


def _safe(v):
    """Neutralise spreadsheet-formula injection in exported strings."""
    if isinstance(v, str) and v[:1] in "=+-@":
        return "'" + v
    return "" if v is None else v


def _pct(v):
    """Delegates to the single canonical Python formatter (sla.fmt_pct), which mirrors the
    UI's frontend/src/format.ts. Kept as a thin alias so report code reads naturally."""
    return sla.fmt_pct(v)


def _status_hex(text: str) -> str:
    """Map a status/trend word to a palette colour (Excel hex)."""
    t = str(text).lower()
    if any(w in t for w in ("target met", "healthy", "improv", "above", "met", "monitored")) and "not " not in t:
        return GREEN
    if any(w in t for w in ("critical", "below", "worsen", "breach", "down", "fault")):
        return RED
    if any(w in t for w in ("risk", "warn", "degrad", "partial", "insufficient")):
        return AMBER
    return SLATE


def _snapshot(db: Session) -> dict:
    """Pull one coherent analytics snapshot so every sheet/slide reconciles."""
    return {"overview": overview.build(db), "sla": sla.fleet_sla(db), "telemetry": telemetry.build(db),
            "trends": trends.build(db), "resilience": resilience.latest_assessment(db)}


def _mlabel(iso):
    try:
        return datetime.fromisoformat(iso).strftime("%b %Y") if iso else None
    except (TypeError, ValueError):
        return None


def _agg_metrics(rows, target):
    """Availability + SLA-budget metrics for a set of daily rows. 'Budget consumed' is the
    share of the ALLOWED downtime that was actually used — the metric that matters for
    SLA-based payments (over 100% means the SLA was breached)."""
    a = sla._aggregate(rows)
    observed, up = a["observed_minutes"], a["up_minutes"]
    # Compliance-adjacent values are withheld with availability. A partial collection
    # must never read as zero downtime or zero budget consumed.
    sufficient = a["availability"] is not None
    down = max(0, observed - up) if sufficient else None
    allowed = observed * (1 - target / 100.0) if sufficient else None
    return {"availability": a["availability"], "coverage": a["coverage"], "observed": observed, "up": up,
            "down": down, "budget_pct": (round(100.0 * down / allowed, 1) if allowed and allowed > 0 else None),
            "first_observed": a.get("first_observed")}


def report_model(db: Session) -> dict:
    """One rich, executive-grade SLA dataset shared by the Excel and PowerPoint reports:
    fleet + business-unit + per-branch WTD/YTD availability, SLA compliance vs target,
    downtime-budget consumed, incidents (per branch and per month), and MTTR — everything a
    non-technical leader needs to compare branches and settle SLA-linked payments."""
    settings = get_settings(); target = settings.sla_target
    ref = sla.today_local()
    ys, ye = sla._window_bounds("ytd", ref); ws, we = sla._window_bounds("wtd", ref)
    d30s, d30e = sla._window_bounds("rolling_30", ref)
    fleet = overview._fleet(db)
    ids = [d["device_id"] for d in fleet if d["device_id"]]
    by_dev = defaultdict(list)
    if ids:
        for r in db.scalars(select(SlaDaily).where(SlaDaily.device_id.in_(ids), SlaDaily.day >= ys, SlaDaily.day <= ye)).all():
            by_dev[r.device_id].append(r)

    def win(dev_ids, s, e):
        return _agg_metrics([r for did in dev_ids for r in by_dev.get(did, []) if s <= r.day <= e], target)

    inc = trends.incidents(db, fleet=fleet)
    device_site = {d["device_id"]: d["site_code"] or d["city"] for d in fleet if d["device_id"]}
    by_site = defaultdict(lambda: {"count": 0, "down": 0})
    for i in inc:
        site = device_site.get(i["device_id"])
        if site:
            by_site[site]["count"] += 1; by_site[site]["down"] += i["down"]

    # per-branch (site) scorecard
    site_groups = defaultdict(list)
    for d in fleet:
        site_groups[d["site_code"] or d["city"]].append(d)
    branches = []
    for members in site_groups.values():
        first = members[0]; did = [m["device_id"] for m in members if m["device_id"]]
        wtd, ytd = win(did, ws, we), win(did, ys, ye)
        ci = by_site.get(first["site_code"] or first["city"], {"count": 0, "down": 0})
        branches.append({
            "site_code": first["site_code"], "city": first["city"], "province": first["province"], "unit": first["business_unit"],
            "devices": sum(m["physical"] for m in members), "wtd": wtd["availability"], "ytd": ytd["availability"], "target": target,
            "met_wtd": (wtd["availability"] >= target) if wtd["availability"] is not None else None,
            "met_ytd": (ytd["availability"] >= target) if ytd["availability"] is not None else None,
            "down_ytd_min": ytd["down"], "budget_pct": ytd["budget_pct"], "incidents": ci["count"],
            "mttr": (round(ci["down"] / ci["count"]) if ci["count"] else None), "since": _mlabel(ytd["first_observed"]),
        })
    branches.sort(key=lambda b: (b["ytd"] is None, -(b["ytd"] or 0)))
    for i, b in enumerate(branches, 1):
        b["rank"] = i

    # business units (Healthcare vs Dental)
    unit_groups = defaultdict(list)
    for d in fleet:
        unit_groups[d["business_unit"] or "Unassigned"].append(d)
    units = []
    for name, members in unit_groups.items():
        did = [m["device_id"] for m in members if m["device_id"]]
        wtd, ytd = win(did, ws, we), win(did, ys, ye)
        sites = {m["site_code"] or m["city"] for m in members}
        ui = sum(by_site.get(s, {"count": 0})["count"] for s in sites)
        ud = sum(by_site.get(s, {"down": 0})["down"] for s in sites)
        units.append({"name": name, "wtd": wtd["availability"], "ytd": ytd["availability"], "devices": sum(m["physical"] for m in members),
                      "sites": len({m["site_code"] or m["city"] for m in members}), "incidents": ui,
                      "mttr": (round(ud / ui) if ui else None), "budget_pct": ytd["budget_pct"],
                      "met_ytd": (ytd["availability"] >= target) if ytd["availability"] is not None else None})
    units.sort(key=lambda u: -u["devices"])

    # Incidents per calendar month (last 6); zero months are real chart data, not gaps.
    per_month = defaultdict(lambda: {"count": 0, "down": 0})
    for i in inc:
        per_month[i["start"][:7]]["count"] += 1; per_month[i["start"][:7]]["down"] += i["down"]
    months = []
    for offset in range(5, -1, -1):
        year, month = ref.year, ref.month - offset
        while month <= 0:
            year -= 1; month += 12
        months.append(f"{year:04d}-{month:02d}")
    inc_month = [{"month": m, "count": per_month[m]["count"], "down": per_month[m]["down"]} for m in months]

    f_wtd, f_ytd = win(ids, ws, we), win(ids, ys, ye)
    crit_ids = [d["device_id"] for d in fleet if d["band"] == "Critical" and d["device_id"]]
    return {
        "target": target, "stamp": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M"),
        "ytd_start": ys.isoformat(), "wtd_start": ws.isoformat(), "as_of": ref.isoformat(),
        "fleet": {"wtd": f_wtd["availability"], "ytd": f_ytd["availability"], "budget_pct": f_ytd["budget_pct"],
                   "down_ytd_min": f_ytd["down"], "met_wtd": f_wtd["availability"] is not None and f_wtd["availability"] >= target,
                   "met_ytd": f_ytd["availability"] is not None and f_ytd["availability"] >= target},
        "units": units, "branches": branches, "incidents_month": inc_month,
        "weekly": trends.availability_trend(db, fleet=fleet)["series"],
        "mttr_mtbf": trends.mttr_mtbf(db, inc=inc, fleet=fleet), "deltas": trends.deltas(db, fleet=fleet),
        "branches_met": sum(1 for b in branches if b["met_ytd"] is True),
        "branches_missed": sum(1 for b in branches if b["met_ytd"] is False),
        "branches_measured": sum(1 for b in branches if b["met_ytd"] is not None),
        "critical": (win(crit_ids, ys, ye) if crit_ids else None), "critical_devices": len(crit_ids),
    }


def _readout(o: dict, tr: dict) -> list:
    """Executive-readout action rows (priority / risk / business impact / recommendation)."""
    rows = []
    for i, a in enumerate(o["actions"][:6], 1):
        rows.append((f"P{a.get('priority', i)}", a["title"], a["severity"], a["detail"]))
    if not rows:
        rows.append(("—", "No leadership actions", "OK", "Network is operating within targets."))
    return rows


# =========================================================================== Excel

_THIN = Side(style="thin", color=LINE)
_BORDER = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)


def _dashboard_sheet(wb, o, tr, m, stamp):
    """The Executive Dashboard cover: brand banner, 6 KPI tiles, and a readout table."""
    ws = wb.create_sheet("Executive Dashboard")
    ws.sheet_view.showGridLines = False
    for col in range(1, 13):
        ws.column_dimensions[get_column_letter(col)].width = 12.5
    g, h = o["global_sla"], o["header"]

    # brand banner
    ws.merge_cells("A1:C4")
    ws["A1"] = "MEDLINE"; ws["A1"].fill = PatternFill("solid", fgColor=BRAND)
    ws["A1"].font = Font(size=20, bold=True, color="FFFFFF"); ws["A1"].alignment = Alignment("center", "center")
    ws.merge_cells("D1:L2"); ws.merge_cells("D3:L4")
    ws["D1"] = f"{h['org']} — Enterprise Network Reliability"
    ws["D1"].font = Font(size=22, bold=True, color="FFFFFF")
    ws["D3"] = f"Executive SLA dashboard · generated {stamp} UTC · data quality {h['data_quality']}"
    ws["D3"].font = Font(size=11, color="D9EAFB")
    for cell in ("D1", "D3"):
        ws[cell].alignment = Alignment("left", "center")
    for col in range(1, 13):
        for r in (1, 2, 3, 4):
            ws.cell(row=r, column=col).fill = PatternFill("solid", fgColor=BRAND)
    for r in (1, 2, 3, 4):
        ws.row_dimensions[r].height = 20

    # KPI tiles (6, two columns each)
    fl = m["fleet"]; bud = fl["budget_pct"]
    tiles = [
        ("Network SLA — This Week", _pct(fl["wtd"]), "Target met" if fl["met_wtd"] else "Below target"),
        ("Network SLA — Year to Date", _pct(fl["ytd"]), "Target met" if fl["met_ytd"] else "Below target"),
        ("SLA Target", f"{m['target']}%", "Objective"),
        ("SLA Budget Used", (f"{bud}%" if bud is not None else "—"), "Within budget" if (bud is not None and bud <= 100) else "Over budget"),
        ("Branches Meeting SLA", f"{m['branches_met']}/{len(m['branches'])}", f"{m['branches_missed']} missed"),
        ("Incidents (90 days)", str(m["mttr_mtbf"]["incidents"]), f"MTTR {m['mttr_mtbf']['mttr_minutes']} min"),
    ]
    starts = ["A", "C", "E", "G", "I", "K"]
    for (label, value, status), c in zip(tiles, starts):
        c2 = get_column_letter(ord(c) - 64 + 1)  # tile spans 2 columns
        ws.merge_cells(f"{c}6:{c2}6"); ws.merge_cells(f"{c}7:{c2}8"); ws.merge_cells(f"{c}9:{c2}9")
        ws[f"{c}6"] = label.upper(); ws[f"{c}6"].font = Font(size=8, bold=True, color=SLATE)
        ws[f"{c}7"] = value; ws[f"{c}7"].font = Font(size=20, bold=True, color=NAVY)
        ws[f"{c}7"].alignment = Alignment("left", "center")
        ws[f"{c}9"] = status; ws[f"{c}9"].font = Font(size=8, bold=True, color=_status_hex(status))
        for rr in (6, 7, 8, 9):
            for cc in range(ord(c) - 64, ord(c2) - 64 + 1):
                cell = ws.cell(row=rr, column=cc)
                cell.fill = PatternFill("solid", fgColor="FFFFFF")
                cell.border = _BORDER
    ws.row_dimensions[7].height = 24

    # Executive readout narrative + action table
    ws["A11"] = "Executive Readout"; ws["A11"].font = Font(size=14, bold=True, color=NAVY)
    ws.merge_cells("A12:E22")
    ws["A12"] = o["summary"]; ws["A12"].alignment = Alignment("left", "top", wrap_text=True)
    ws["A12"].font = Font(size=11, color=NAVY)
    headers = ["Priority", "Risk", "Severity", "Recommended action"]
    for i, hdr in enumerate(headers):
        cell = ws.cell(row=12, column=7 + i, value=hdr)
        cell.font = Font(size=9, bold=True, color=SLATE); cell.fill = PatternFill("solid", fgColor=ZEBRA)
    for r, (pri, risk, sev, act) in enumerate(_readout(o, tr), 13):
        for i, val in enumerate((pri, risk, sev, act)):
            cell = ws.cell(row=r, column=7 + i, value=_safe(val))
            cell.alignment = Alignment("left", "top", wrap_text=True)
            cell.font = Font(size=9, color=_status_hex(sev) if i == 2 else NAVY, bold=(i == 2))
    for col, w in (("G", 8), ("H", 26), ("I", 12), ("J", 40)):
        ws.column_dimensions[col].width = w
    return ws


def _sheet(wb, title, subtitle=""):
    ws = wb.create_sheet(title[:31])
    ws.sheet_view.showGridLines = False
    ws.merge_cells("A1:F1")
    ws["A1"] = title; ws["A1"].font = Font(size=15, bold=True, color="FFFFFF")
    for col in range(1, 7):
        ws.cell(row=1, column=col).fill = PatternFill("solid", fgColor=BRAND)
    ws.row_dimensions[1].height = 24
    if subtitle:
        ws["A2"] = subtitle; ws["A2"].font = Font(size=10, italic=True, color=SLATE)
    return ws


def _table(ws, headers, rows, row=4):
    for c, hh in enumerate(headers, 1):
        cell = ws.cell(row=row, column=c, value=hh)
        cell.font = Font(size=10, bold=True, color="FFFFFF"); cell.fill = PatternFill("solid", fgColor=SLATE)
    zebra = PatternFill("solid", fgColor=ZEBRA)
    for r, data in enumerate(rows, row + 1):
        for c, value in enumerate(data, 1):
            cell = ws.cell(row=r, column=c, value=_safe(value))
            if (r - row) % 2 == 0:
                cell.fill = zebra
    for i in range(len(headers)):
        ws.column_dimensions[get_column_letter(i + 1)].width = 22
    ws.freeze_panes = ws.cell(row=row + 1, column=1)
    if rows:
        ws.auto_filter.ref = f"A{row}:{get_column_letter(len(headers))}{row + len(rows)}"


def _met(v):
    return "Yes" if v is True else ("No" if v is False else "Insufficient")


def _colored_table(ws, headers, rows, colorers, row=4):
    """Like _table but a `colorers` dict {col_index: fn(value)->hex or None} tints cells."""
    for c, hh in enumerate(headers, 1):
        cell = ws.cell(row=row, column=c, value=hh)
        cell.font = Font(size=10, bold=True, color="FFFFFF"); cell.fill = PatternFill("solid", fgColor=SLATE)
    zebra = PatternFill("solid", fgColor=ZEBRA)
    for r, data in enumerate(rows, row + 1):
        for c, value in enumerate(data, 1):
            cell = ws.cell(row=r, column=c, value=_safe(value))
            if (r - row) % 2 == 0:
                cell.fill = zebra
            clr = colorers.get(c - 1, lambda v: None)(value) if colorers else None
            if clr:
                cell.font = Font(bold=True, color=clr)
    for i in range(len(headers)):
        ws.column_dimensions[get_column_letter(i + 1)].width = 16 if i else 8
    ws.freeze_panes = ws.cell(row=row + 1, column=1)
    if rows:
        ws.auto_filter.ref = f"A{row}:{get_column_letter(len(headers))}{row + len(rows)}"


def executive_excel(db: Session) -> Path:
    d = _snapshot(db)
    o, s, tel, tr = d["overview"], d["sla"], d["telemetry"], d["trends"]
    m = report_model(db)
    stamp = datetime.now(timezone.utc).strftime("%Y-%m-%d_%H%M%S")
    folder = Path(get_settings().report_dir); folder.mkdir(parents=True, exist_ok=True)
    path = folder / f"Executive_Reliability_{stamp}.xlsx"
    wb = Workbook(); wb.remove(wb.active)
    met_color = lambda v: (GREEN if v == "Yes" else RED)

    _dashboard_sheet(wb, o, tr, m, stamp)

    # ---- Branch SLA Scorecard — the payment-relevant leaderboard ----
    ws = _sheet(wb, "Branch SLA Scorecard", "Every branch ranked by year-to-date availability. 'Met SLA' compares against the target. 'Budget Used' is the share of the allowed downtime consumed (over 100% = breach).")
    _colored_table(ws, ["Rank", "Branch", "City", "Business Unit", "This Week", "Year to Date", "Target", "Met SLA (WTD)", "Met SLA (YTD)", "Downtime YTD (min)", "SLA Budget Used", "Incidents", "Avg Resolve (min)"],
                   [(b["rank"], b["site_code"], b["city"] + (f" · since {b['since']}" if b.get("since") else ""), b["unit"],
                     _pct(b["wtd"]), _pct(b["ytd"]), f"{b['target']}%", _met(b["met_wtd"]), _met(b["met_ytd"]),
                     b["down_ytd_min"], (f"{b['budget_pct']}%" if b["budget_pct"] is not None else "—"), b["incidents"],
                     (b["mttr"] if b["mttr"] is not None else "—")) for b in m["branches"]],
                   {7: met_color, 8: met_color})

    # ---- Healthcare vs Dental ----
    ws = _sheet(wb, "Healthcare vs Dental", "How the two business units compare on reliability, SLA budget and incidents.")
    _colored_table(ws, ["Business Unit", "This Week", "Year to Date", "Target", "Met SLA (YTD)", "SLA Budget Used", "Devices", "Sites", "Incidents (90d)", "Avg Resolve (min)"],
                   [(u["name"], _pct(u["wtd"]), _pct(u["ytd"]), f"{m['target']}%", _met(u["met_ytd"]),
                     (f"{u['budget_pct']}%" if u["budget_pct"] is not None else "—"), u["devices"], u["sites"], u["incidents"],
                     (u["mttr"] if u["mttr"] is not None else "—")) for u in m["units"]],
                   {4: met_color})

    # ---- Weekly Trend ----
    ws = _sheet(wb, "Weekly Trend", "Fleet availability and downtime, week by week (last 12 weeks).")
    _table(ws, ["Week Starting", "Availability", "Vs Target", "Downtime (min)"],
           [(w["week_start"], _pct(w["availability"]), ("Below" if w["below_target"] else "Met"), w["downtime_minutes"]) for w in m["weekly"]])

    # ---- Incidents by Month ----
    ws = _sheet(wb, "Incidents by Month", "How many incidents were logged each month, and total downtime.")
    _table(ws, ["Month", "Incidents", "Total Downtime (min)"], [(x["month"], x["count"], x["down"]) for x in m["incidents_month"]] or [("No incidents in range", 0, 0)])

    # ---- Supporting detail ----
    ws = _sheet(wb, "SLA by Device", "Per-device WTD/YTD availability, coverage-gated")
    _table(ws, ["Device", "Site", "WTD", "YTD", "Coverage YTD"],
           [(e["hostname"], e["site"], _pct(e["wtd"]["availability"]), _pct(e["ytd"]["availability"]), _pct(e["ytd"]["coverage"])) for e in s["devices"]])

    ws = _sheet(wb, "Critical Applications", "Business-application SLAs (mapping pending). Critical-criticality infrastructure is shown as today's closest measured proxy.")
    crit = m["critical"]
    _table(ws, ["Item", "Measurement", "This Week / YTD", "Note"], [
        ("Business application SLIs", "Per-application availability", "Mapping pending", "Populates when an authoritative LogicMonitor application SLI is mapped — never inferred from device health."),
        (f"Critical infrastructure ({m['critical_devices']} devices)", "Availability YTD (proxy)", _pct(crit["availability"]) if crit else "Insufficient", "The most business-critical devices, as the closest real measurement available today."),
    ])

    ws = _sheet(wb, "Data Coverage", "What LogicMonitor exposes for this fleet")
    _table(ws, ["Metric", "LogicMonitor source", "Coverage", "Devices"],
           [(x["metric"], x["source"], x["status"], f"{x['devices']}/{tel['data_coverage']['fleet']}") for x in tel["data_coverage"]["metrics"]])

    ws = _sheet(wb, "SLA Methodology", "How every number is measured — the basis for SLA-linked payments.")
    _table(ws, ["Topic", "How it is measured"], [
        ("Availability", "Up eligible minutes ÷ observed eligible minutes × 100. A minute is 'up' when the device answered at least one probe (packet loss < 100%)."),
        ("This Week (WTD)", "Monday to now, in the configured timezone (America/Vancouver)."),
        ("Year to Date (YTD)", "From each branch's first monitored day to now — a mid-year branch is not penalised for months it did not exist."),
        ("Coverage gating", "If less than 90% of the expected minutes were observed, no number is published (shown as Insufficient) — never a fabricated 0% or 100%."),
        ("SLA Budget Used", "Downtime ÷ allowed downtime × 100, where allowed downtime = observed minutes × (1 − target). Over 100% means the SLA was breached."),
        ("Incidents / Avg Resolve", "Incidents are runs of below-100% availability days; Avg Resolve (MTTR) is total downtime ÷ number of incidents. Availability-derived and approximate."),
        ("Scope", "Only devices registered in Device Fleet are analysed. WAN provider routers are excluded from Medline SLA."),
        ("Target", f"The SLA target is {m['target']}% (configurable)."),
    ])
    wb.save(path)
    return path


# ====================================================================== PowerPoint

def _rect(slide, l, t, w, h, fill, line=None, rounded=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                                   Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _text(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = align
        r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return box


def _kpi(slide, l, t, w, label, value, sub, value_color=_NAVY):
    _rect(slide, l, t, w, 1.15, _WHITE, line=_LINE)
    _text(slide, l + 0.12, t + 0.10, w - 0.24, 0.25, label.upper(), 9, _SLATE, bold=True)
    _text(slide, l + 0.12, t + 0.34, w - 0.24, 0.5, value, 24, value_color, bold=True)
    _text(slide, l + 0.12, t + 0.84, w - 0.24, 0.25, sub, 8.5, _GREY)


def _chip(slide, l, t, text, fill_hex, text_rgb):
    w = 0.16 + 0.075 * len(text)
    _rect(slide, l, t, w, 0.26, RGBColor.from_string(fill_hex))
    _text(slide, l, t + 0.02, w, 0.22, text, 8.5, text_rgb, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return l + w + 0.12


def _header(slide, title, subtitle, page):
    _rect(slide, 0, 0, 13.333, 7.5, _CANVAS, rounded=False)
    _rect(slide, 0, 0, 13.333, 0.16, _BRAND, rounded=False)
    _rect(slide, 0, 0, 0.6, 0.6, _BRAND, rounded=False)
    _text(slide, 0.8, 0.18, 10, 0.4, title, 22, _NAVY, bold=True)
    _text(slide, 0.8, 0.66, 11.5, 0.3, subtitle, 12, _GREY)
    _text(slide, 12.4, 7.05, 0.5, 0.3, page, 10, _GREY, align=PP_ALIGN.RIGHT)


def _bar_chart(slide, l, t, w, h, sites, target):
    """Per-site availability bars, scaled within a tight window so differences read; below
    target is red. Draws on top of a white card the caller already placed."""
    usable = [x for x in sites if isinstance(x.get("availability_ytd"), (int, float))]
    usable = sorted(usable, key=lambda x: x["availability_ytd"])[:8]
    if not usable:
        _text(slide, l, t + h / 2, w, 0.3, "Insufficient evidence to chart", 11, _GREY, align=PP_ALIGN.CENTER)
        return
    lo = min(min(x["availability_ytd"] for x in usable), target) - 0.05
    hi = 100.0
    span = max(hi - lo, 0.01)
    n = len(usable)
    gap = 0.18
    bw = (w - gap * (n + 1)) / n
    base = t + h - 0.35
    # target line
    ty = base - (h - 0.5) * (target - lo) / span
    _rect(slide, l, ty, w, 0.02, _AMBER, rounded=False)
    for i, x in enumerate(usable):
        val = x["availability_ytd"]
        bh = max(0.1, (h - 0.5) * (val - lo) / span)
        bx = l + gap + i * (bw + gap)
        color = _RED if val < target else _BRAND
        _rect(slide, bx, base - bh, bw, bh, color, rounded=False)
        _text(slide, bx - 0.1, base + 0.02, bw + 0.2, 0.25, x.get("site_code") or x.get("city", ""), 8, _SLATE, align=PP_ALIGN.CENTER)


def _vbars(slide, l, t, w, h, data, lo, hi, target=None, valfmt=None):
    """Simple vertical bar chart. data = list of (label, value_or_None, color_rgb)."""
    n = max(len(data), 1); gap = 0.08; bw = (w - gap * (n + 1)) / n; base = t + h - 0.34; span = max(hi - lo, 0.001)
    if target is not None:
        ty = base - (h - 0.5) * (target - lo) / span
        _rect(slide, l, ty, w, 0.02, _AMBER, rounded=False)
    for i, (lab, val, col) in enumerate(data):
        bx = l + gap + i * (bw + gap)
        if isinstance(val, (int, float)):
            bh = max(0.08, (h - 0.5) * (val - lo) / span)
            _rect(slide, bx, base - bh, bw, bh, col, rounded=False)
            if valfmt:
                _text(slide, bx - 0.15, base - bh - 0.24, bw + 0.3, 0.22, valfmt(val), 7.5, _SLATE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, bx - 0.15, base + 0.03, bw + 0.3, 0.22, lab, 7.5, _SLATE, align=PP_ALIGN.CENTER)


def executive_pptx(db: Session) -> Path:
    d = _snapshot(db)
    o, tel, tr = d["overview"], d["telemetry"], d["trends"]
    h = o["header"]
    m = report_model(db); fl = m["fleet"]; tgt = m["target"]
    stamp = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    folder = Path(get_settings().report_dir); folder.mkdir(parents=True, exist_ok=True)
    path = folder / f"Executive_Reliability_{datetime.now(timezone.utc).strftime('%Y-%m-%d_%H%M%S')}.pptx"
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    metchip = lambda ok: (GREEN_BG, _GREEN, "MET") if ok else (RED_BG, _RED, "MISSED")

    # ---- Slide 1: cover ----
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, _BRAND, rounded=False)
    logo = Path(__file__).parent / "medline-logo.jpg"
    if logo.exists():
        s.shapes.add_picture(str(logo), Inches(0.8), Inches(0.7), height=Inches(1.9))
    _text(s, 5.2, 2.2, 7.4, 1.6, "Network SLA\nExecutive Report", 38, _WHITE, bold=True)
    _text(s, 5.25, 4.05, 8, 0.35, f"{h['org']}   |   Reliability & Service-Level Performance", 14, RGBColor(0xD9, 0xEA, 0xFB), bold=True)
    _text(s, 5.25, 4.5, 8, 0.3, f"Generated {stamp} UTC  ·  Week-to-date and Year-to-date", 12, RGBColor(0xBE, 0xD3, 0xF2))

    # ---- Slide 2: executive scorecard ----
    s = prs.slides.add_slide(blank)
    _header(s, "Executive Scorecard", "The 30-second read: are we meeting our service levels this week and this year?", "02")
    tiles = [("Network SLA — This Week", _pct(fl["wtd"]), "target " + f"{tgt}%", _GREEN if fl["met_wtd"] else _RED),
             ("Network SLA — Year to Date", _pct(fl["ytd"]), "target " + f"{tgt}%", _GREEN if fl["met_ytd"] else _RED),
             ("SLA Budget Used", (f"{fl['budget_pct']}%" if fl["budget_pct"] is not None else "—"), "of allowed downtime", _NAVY if (fl["budget_pct"] or 0) <= 100 else _RED),
             ("Branches Meeting SLA", f"{m['branches_met']}/{len(m['branches'])}", f"{m['branches_missed']} need attention", _NAVY),
             ("Incidents (90 days)", str(m["mttr_mtbf"]["incidents"]), f"avg resolve {m['mttr_mtbf']['mttr_minutes']} min", _NAVY),
             ("Trend", m["deltas"]["wow"]["trend"].title(), "week over week", _GREEN if "improv" in m["deltas"]["wow"]["trend"].lower() else (_RED if "worsen" in m["deltas"]["wow"]["trend"].lower() else _NAVY))]
    lx = 0.8
    for lab, val, sub, col in tiles:
        _kpi(s, lx, 1.2, 1.95, lab, val, sub, col); lx += 2.05
    _rect(s, 0.8, 2.65, 5.2, 3.9, _WHITE, line=_LINE)
    _text(s, 1.1, 2.85, 4.6, 0.3, "What this means", 14, _NAVY, bold=True)
    _text(s, 1.1, 3.3, 4.7, 2.2, o["summary"], 11, _SLATE)
    cx = _chip(s, 1.1, 6.05, "This week " + ("MET" if fl["met_wtd"] else "MISSED"), *( (GREEN_BG,_GREEN) if fl["met_wtd"] else (RED_BG,_RED)))
    _chip(s, cx, 6.05, "YTD " + ("MET" if fl["met_ytd"] else "MISSED"), *((GREEN_BG,_GREEN) if fl["met_ytd"] else (RED_BG,_RED)))
    _rect(s, 6.3, 2.65, 6.2, 3.9, _WHITE, line=_LINE)
    _text(s, 6.6, 2.85, 5.6, 0.3, "Branch availability vs. SLA target (YTD)", 13, _NAVY, bold=True)
    bd = [(b["site_code"] or b["city"][:4], b["ytd"], (_BRAND if b["met_ytd"] else _RED)) for b in sorted(m["branches"], key=lambda x:(x["ytd"] is None,-(x["ytd"] or 0)))[:12]]
    lo = min([v for _,v,_ in bd if isinstance(v,(int,float))] + [tgt]) - 0.05
    _vbars(s, 6.55, 3.35, 5.7, 3.0, bd, lo, 100.0, target=tgt)

    # ---- Slide 3: branch SLA leaderboard ----
    s = prs.slides.add_slide(blank)
    _header(s, "Branch SLA Leaderboard", "Every branch, ranked. Green met the target; red missed it. This is the basis for SLA-linked measures.", "03")
    ty = 1.4
    cols = [(0.9, "Rank"), (1.7, "Branch"), (3.4, "This Week"), (5.1, "Year to Date"), (7.0, "Met YTD"), (8.6, "Incidents"), (10.4, "Avg Resolve")]
    _rect(s, 0.8, ty, 11.9, 0.34, _BRAND)
    for cxx, head in cols:
        _text(s, cxx, ty + 0.05, 2, 0.26, head, 10, _WHITE, bold=True)
    for i, b in enumerate(m["branches"][:14]):
        yy = ty + 0.42 + i * 0.34
        if i % 2 == 0:
            _rect(s, 0.8, yy - 0.03, 11.9, 0.34, RGBColor.from_string(ZEBRA))
        cells = [str(b["rank"]), f"{b['site_code']} · {b['city']}", _pct(b["wtd"]), _pct(b["ytd"]),
                 "MET" if b["met_ytd"] else "MISSED", str(b["incidents"]), (f"{b['mttr']} min" if b["mttr"] is not None else "—")]
        for (cxx, _), val, idx in zip(cols, cells, range(len(cells))):
            col = (_GREEN if b["met_ytd"] else _RED) if idx == 4 else _NAVY
            _text(s, cxx, yy, 2.0, 0.3, val, 9.5, col, bold=(idx == 4))

    # ---- Slide 4: Healthcare vs Dental ----
    s = prs.slides.add_slide(blank)
    _header(s, "Healthcare vs Dental", "How the two business units compare — reliability, SLA budget, and how quickly incidents are resolved.", "04")
    lx = 0.8
    for u in m["units"][:2]:
        _rect(s, lx, 1.3, 5.9, 4.0, _WHITE, line=_LINE)
        _text(s, lx + 0.3, 1.5, 5.3, 0.35, u["name"].upper(), 15, _BRAND, bold=True)
        _text(s, lx + 0.3, 1.95, 3, 0.6, _pct(u["ytd"]), 34, _NAVY, bold=True)
        _text(s, lx + 0.3, 2.65, 5, 0.3, "Year-to-date availability", 10, _GREY)
        bg, fg, lbl = metchip(u["met_ytd"]); _chip(s, lx + 3.7, 2.15, "SLA " + lbl, bg, fg)
        stats = [("This week", _pct(u["wtd"])), ("SLA budget used", (f"{u['budget_pct']}%" if u["budget_pct"] is not None else "—")),
                 ("Devices / sites", f"{u['devices']} / {u['sites']}"), ("Incidents (90d)", str(u["incidents"])),
                 ("Avg resolve time", (f"{u['mttr']} min" if u["mttr"] is not None else "—"))]
        yy = 3.15
        for k, v in stats:
            _text(s, lx + 0.3, yy, 3.2, 0.3, k, 11, _SLATE)
            _text(s, lx + 3.6, yy, 2.0, 0.3, v, 11, _NAVY, bold=True)
            yy += 0.42
        lx += 6.1
    _text(s, 0.8, 5.55, 11.9, 0.8, "Both units are compared against the same " + f"{tgt}% target. 'SLA budget used' is how much of the allowed downtime each unit has consumed — the closer to 100%, the closer to a breach.", 11, _SLATE)

    # ---- Slide 5: trend & incidents ----
    s = prs.slides.add_slide(blank)
    _header(s, "Trend & Incidents", "Are we getting better or worse? Weekly availability and monthly incident volume.", "05")
    _rect(s, 0.8, 1.35, 6.0, 5.0, _WHITE, line=_LINE)
    _text(s, 1.1, 1.55, 5.4, 0.3, "Availability by week (last 12 weeks)", 13, _NAVY, bold=True)
    wk = [(w["week_start"][5:], w["availability"], (_RED if w["below_target"] else _BRAND)) for w in m["weekly"]]
    wlo = min([v for _,v,_ in wk if isinstance(v,(int,float))] + [tgt]) - 0.05
    _vbars(s, 1.0, 2.05, 5.6, 4.0, wk, wlo, 100.0, target=tgt)
    _rect(s, 7.0, 1.35, 5.5, 5.0, _WHITE, line=_LINE)
    _text(s, 7.3, 1.55, 4.9, 0.3, "Incidents logged per month", 13, _NAVY, bold=True)
    im = m["incidents_month"]
    mx = max([x["count"] for x in im] + [1])
    mb = [(x["month"][5:], x["count"], _SLATE) for x in im]
    _vbars(s, 7.2, 2.05, 5.1, 4.0, mb, 0, mx * 1.15, valfmt=lambda v: str(int(v)))
    if len(im) >= 2 and im[-1]["count"] <= im[0]["count"]:
        _text(s, 7.3, 6.0, 5.0, 0.4, f"Trend: incidents fell from {im[0]['count']} in {im[0]['month']} to {im[-1]['count']} in {im[-1]['month']}.", 10, _GREEN, bold=True)

    # ---- Slide 6: the metrics that matter ----
    s = prs.slides.add_slide(blank)
    _header(s, "The metrics that matter", "Good SLA reporting is more than weekly uptime. These measures tell the real story.", "06")
    cards = [
        ("Availability (WTD / YTD)", f"{_pct(fl['wtd'])} / {_pct(fl['ytd'])}", "Was the service working? The headline number."),
        ("SLA compliance", f"{m['branches_met']}/{len(m['branches'])} branches", "How many locations met the agreed target."),
        ("SLA budget used", (f"{fl['budget_pct']}%" if fl['budget_pct'] is not None else "—"), "How much of the allowed downtime we've spent (100% = breach)."),
        ("Avg time to resolve", f"{m['mttr_mtbf']['mttr_minutes']} min", "How fast we recover when something breaks (MTTR)."),
        ("Incident frequency", f"{m['mttr_mtbf']['incidents']} in 90 days", "How often things break — falling is good (MTBF)."),
        ("Consistency", f"Best {m['branches'][0]['city']} · Worst {m['branches'][-1]['city']}", "Steady beats a good average with one big outage."),
    ]
    for i, (lab, val, sub) in enumerate(cards):
        cx = 0.8 + (i % 3) * 4.05; cy = 1.5 + (i // 3) * 2.5
        _rect(s, cx, cy, 3.85, 2.2, _WHITE, line=_LINE)
        _text(s, cx + 0.2, cy + 0.15, 3.5, 0.3, lab.upper(), 10, _SLATE, bold=True)
        _text(s, cx + 0.2, cy + 0.5, 3.5, 0.6, val, 20, _NAVY, bold=True)
        _text(s, cx + 0.2, cy + 1.25, 3.5, 0.8, sub, 10.5, _GREY)

    # ---- Slide 7: methodology / payment basis ----
    s = prs.slides.add_slide(blank)
    _header(s, "How every number is measured", "The basis for SLA reporting and any SLA-linked measures.", "07")
    defs = [
        ("Availability", "Up minutes ÷ observed minutes × 100. A minute counts as 'up' when the device answered at least one probe."),
        ("This Week / Year to Date", "This Week is Monday to now. Year to Date runs from each branch's first monitored day — a new branch is not penalised for months it did not exist."),
        ("SLA budget used", "Downtime ÷ allowed downtime × 100. Allowed = observed minutes × (1 − target). Over 100% means the SLA was breached."),
        ("Evidence gating", "If less than 90% of the expected minutes were observed, no number is published — never a fabricated 0% or 100%."),
        ("Incidents & resolve time", "Incidents are runs of below-100% availability days; average resolve time (MTTR) is total downtime ÷ incidents. Availability-derived and approximate."),
        ("Scope", f"Only Medline's own devices are measured, against a {tgt}% target. WAN provider routers are excluded from Medline SLA."),
    ]
    yy = 1.4
    for term, defn in defs:
        _rect(s, 0.8, yy, 11.9, 0.82, _WHITE, line=_LINE)
        _text(s, 1.0, yy + 0.12, 3.2, 0.6, term, 12, _BRAND, bold=True)
        _text(s, 4.3, yy + 0.12, 8.2, 0.6, defn, 10.5, _SLATE, anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.9

    prs.save(path)
    return path
